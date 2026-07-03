# backend/services/parser.py
# VERSI SEDERHANA UNTUK PEMULA

import hashlib
import re
from pathlib import Path
from typing import List, Dict, Optional
from dataclasses import dataclass, field
from datetime import datetime
import logging
import docx
from pptx import Presentation
import openpyxl
import pytesseract
from PIL import Image
import cv2

import fitz  # PyMuPDF

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

@dataclass
class PageData:
    """Data per halaman"""
    page_num: int
    text: str
    has_ocr: bool = False

@dataclass
class Chunk:
    """Potongan teks"""
    text: str
    page_num: int
    chunk_index: int

@dataclass
class ParsedDocument:
    """Hasil parsing"""
    file_path: str
    file_hash: str
    pages: List[PageData]
    chunks: List[Chunk]
    total_pages: int
    parsing_time: float

class DocumentParser:
    """Parser sederhana untuk PDF"""
    
    def __init__(self):
        self.cache = {}  # Cache sederhana

    def _parse_image(self, filepath: Path) -> ParsedDocument:
        """Parse gambar dengan OCR"""
        # Baca gambar
        image = Image.open(filepath)
    
        # OCR
        text = pytesseract.image_to_string(image, lang='ind')
    
        pages = [PageData(page_num=1, text=text)]
        chunks = self._create_chunks(pages)
    
        return ParsedDocument(
            file_path=str(filepath),
            file_hash=self._compute_hash(filepath),
            pages=pages,
            chunks=chunks,
            total_pages=1,
            parsing_time=0
        )
        
    def parse(self, filepath: str):
        """Parse satu file PDF"""
        start_time = datetime.now()
        filepath = Path(filepath)
        
        # 1. Hitung hash untuk cache
        file_hash = self._compute_hash(filepath)
        cache_key = f"{file_hash}_{filepath.stat().st_size}"
        
        if cache_key in self.cache:
            logger.info(f"✅ Pakai cache: {filepath.name}")
            return self.cache[cache_key]
        
        # 2. Buka PDF
        doc = fitz.open(filepath)
        pages = []
        full_text = ""
        
        # 3. Baca setiap halaman
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            pages.append(PageData(
                page_num=page_num + 1,
                text=text.strip()
            ))
            full_text += text + "\n"
        
        doc.close()
        
        # 4. Buat chunks (potongan teks)
        chunks = self._create_chunks(pages)
        
        # 5. Hasil
        result = ParsedDocument(
            file_path=str(filepath),
            file_hash=file_hash,
            pages=pages,
            chunks=chunks,
            total_pages=len(pages),
            parsing_time=(datetime.now() - start_time).total_seconds()
        )
        
        # 6. Simpan ke cache
        self.cache[cache_key] = result
        
        logger.info(f"✅ Selesai parse {filepath.name}: {len(pages)} halaman, {len(chunks)} chunk")
        return result
    
    def parse_any(self, filepath: str) -> ParsedDocument:
        """Parse berbagai format file"""
        filepath = Path(filepath)
        ext = filepath.suffix.lower()
    
        if ext == '.pdf':
            return self.parse(filepath)
        elif ext == '.docx':
            return self._parse_docx(filepath)
        elif ext == '.pptx':
            return self._parse_pptx(filepath)
        elif ext == '.xlsx':
            return self._parse_xlsx(filepath)
        elif ext == '.txt':
            return self._parse_txt(filepath)
        else:
            raise ValueError(f"Format tidak didukung: {ext}")

def _parse_docx(self, filepath: Path) -> ParsedDocument:
    """Parse DOCX file"""
    doc = docx.Document(filepath)
    
    # Ekstrak teks
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)
    
    text = "\n".join(full_text)
    
    # Buat halaman (simulasi)
    pages = [PageData(page_num=1, text=text)]
    chunks = self._create_chunks(pages)
    
    return ParsedDocument(
        file_path=str(filepath),
        file_hash=self._compute_hash(filepath),
        pages=pages,
        chunks=chunks,
        total_pages=1,
        parsing_time=0
    )

def _parse_pptx(self, filepath: Path) -> ParsedDocument:
    """Parse PPTX file"""
    prs = Presentation(filepath)
    
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                full_text.append(shape.text)
    
    text = "\n".join(full_text)
    pages = [PageData(page_num=1, text=text)]
    chunks = self._create_chunks(pages)
    
    return ParsedDocument(
        file_path=str(filepath),
        file_hash=self._compute_hash(filepath),
        pages=pages,
        chunks=chunks,
        total_pages=1,
        parsing_time=0
    )

def _parse_xlsx(self, filepath: Path) -> ParsedDocument:
    """Parse XLSX file"""
    wb = openpyxl.load_workbook(filepath, data_only=True)
    
    full_text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values=True):
            row_text = [str(cell) for cell in row if cell]
            if row_text:
                full_text.append(" | ".join(row_text))
    
    text = "\n".join(full_text)
    pages = [PageData(page_num=1, text=text)]
    chunks = self._create_chunks(pages)
    
    return ParsedDocument(
        file_path=str(filepath),
        file_hash=self._compute_hash(filepath),
        pages=pages,
        chunks=chunks,
        total_pages=1,
        parsing_time=0
    )

def _parse_txt(self, filepath: Path) -> ParsedDocument:
    """Parse TXT file"""
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
    
    pages = [PageData(page_num=1, text=text)]
    chunks = self._create_chunks(pages)
    
    return ParsedDocument(
        file_path=str(filepath),
        file_hash=self._compute_hash(filepath),
        pages=pages,
        chunks=chunks,
        total_pages=1,
        parsing_time=0
    )
    
    def _create_chunks(self, pages: List[PageData]) -> List[Chunk]:
        """Buat chunks dari teks"""
        chunks = []
        chunk_index = 0
        
        for page in pages:
            # Pecah berdasarkan paragraf
            paragraphs = page.text.split('\n\n')
            
            for para in paragraphs:
                if para.strip():
                    chunks.append(Chunk(
                        text=para.strip(),
                        page_num=page.page_num,
                        chunk_index=chunk_index
                    ))
                    chunk_index += 1
        
        return chunks
    
    def _compute_hash(self, filepath: Path) -> str:
        """Hitung hash file"""
        import hashlib
        hasher = hashlib.sha256()
        with open(filepath, 'rb') as f:
            buf = f.read(65536)
            while len(buf) > 0:
                hasher.update(buf)
                buf = f.read(65536)
        return hasher.hexdigest()[:16]

# ============ FUNGSI UNTUK TESTING ============

def test_parser():
    """Test sederhana"""
    parser = DocumentParser()
    
    # Coba parse file PDF sample
    test_file = "test.pdf"
    
    # Buat sample PDF
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    
    c = canvas.Canvas(test_file, pagesize=letter)
    c.drawString(100, 750, "Surat Dinas")
    c.drawString(100, 700, "Nomor: 440/012/2026")
    c.drawString(100, 675, "Tanggal: 15 Juni 2026")
    c.drawString(100, 650, "Perihal: Cuti Tahunan")
    c.drawString(100, 600, "Kepada Yth. Bapak Budi")
    c.drawString(100, 575, "Dengan hormat, kami mengajukan permohonan cuti...")
    c.save()
    
    print("📄 Testing parser dengan file test.pdf...")
    result = parser.parse(test_file)
    
    print(f"✅ Total halaman: {result.total_pages}")
    print(f"✅ Total chunk: {len(result.chunks)}")
    print(f"✅ Waktu parse: {result.parsing_time:.2f} detik")
    print("\n📝 Chunk pertama:")
    print(result.chunks[0].text[:200] + "...")
    
    # Hapus file test
    import os
    os.remove(test_file)
    
    print("\n✅ Test selesai!")

if __name__ == "__main__":
    test_parser()